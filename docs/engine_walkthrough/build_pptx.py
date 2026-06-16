#!/usr/bin/env python3
"""
build_pptx.py -- assemble the engine-walkthrough figures into an editable .pptx.

No native rasterizer is needed: svglib + reportlab convert each SVG to a high-DPI
PNG (pure Python), and python-pptx places it on a 16:9 slide with an editable
title/subtitle text box. Re-run after editing figures:

    python3 docs/engine_walkthrough/make_figs.py      # refresh the SVGs first
    python3 docs/engine_walkthrough/build_pptx.py     # then rebuild the deck

Writes docs/engine_walkthrough/engine_walkthrough.pptx.
Deps (pip): python-pptx svglib reportlab pillow.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM

HERE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(HERE, "figs")
PNGD = os.path.join(HERE, "_png")
OUT = os.path.join(HERE, "engine_walkthrough.pptx")

INK = RGBColor(0x11, 0x18, 0x27)
MUTE = RGBColor(0x6b, 0x72, 0x80)

# (title, subtitle, figure-svg). Process figures lead; results/companions follow.
SLIDES = [
    ("Stage 0 -- recover the logical schematic from the LPE netlist",
     "de-parasitic R-merge: 689 raw nodes -> 92 logical nets via 1033 R; "
     "164 transistors; bridges = 0 (PASS)",
     "s0_rmerge.svg"),
    ("Union-find -- how raw nodes merge into logical nets",
     "find(x) = which group; union(a,b) = merge; one union per resistor -> "
     "689 raw nodes become 92 logical nets",
     "union_find.svg"),
    ("Channel-Connected Components (CCC) -- group the 92 nets by transistor channels",
     "source-drain channels link nets (gate excluded); rails/inputs are "
     "boundaries; 82 internal nets -> 39 CCCs, 8 with storage",
     "ccc.svg"),
    ("Stage 1 -- HOW storage is found: channel graph (CCC) + feedback loops (SCC)",
     "build the influence digraph, find strongly-connected loops, keep the ones "
     "with >= 2 gate nets = the latches",
     "s1_process.svg"),
    ("Stage 1 -- result: the 8-latch synchronizer chain",
     "8 storage loops ranked by influence-distance to Q -> master ... slave "
     "(4 flip-flops, 2 latches each)",
     "s1_storage.svg"),
    ("Stage 2 -- HOW the bias is derived: switch-level Boolean difference",
     "evaluate the captured node M as inputs toggle -> required vs masked; "
     "no SAT solver",
     "s2_booldiff.svg"),
    ("Stage 2 -- reading the bias vs template.tcl define_arc",
     "the bias is physics-derived, then cross-checked against WHEN; "
     "sensitization is per-arc",
     "s2_sensitize.svg"),
]


def raster(svg_name, dpi=200):
    os.makedirs(PNGD, exist_ok=True)
    png = os.path.join(PNGD, svg_name.replace(".svg", ".png"))
    drawing = svg2rlg(os.path.join(FIGS, svg_name))
    renderPM.drawToFile(drawing, png, fmt="PNG", dpi=dpi)
    return png


def _title_box(slide, text, top, size, color, bold=True):
    tb = slide.shapes.add_textbox(Inches(0.45), top, Inches(12.45), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Arial"
    return tb


def _fit(iw, ih, area_w, area_h):
    """Scale (iw,ih) px to fit (area_w,area_h) EMU preserving aspect; return
    (w_emu, h_emu)."""
    aspect = iw / ih
    if area_w / area_h > aspect:        # height-limited
        h = area_h
        w = int(h * aspect)
    else:                               # width-limited
        w = area_w
        h = int(w / aspect)
    return w, h


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- title slide ----
    s0 = prs.slides.add_slide(blank)
    _title_box(s0, "DeckGen v2 -- how the engine reads a cell BLIND",
               Inches(2.6), 40, INK)
    sub = s0.shapes.add_textbox(Inches(0.45), Inches(3.7), Inches(12.45), Inches(2.0))
    tf = sub.text_frame
    tf.word_wrap = True
    for i, line in enumerate([
        "S0 -> S2 structural walkthrough",
        "cell: SDFSYNC4RPQSXGMZD1BWP130HPNPN3P48CPD   |   corner: "
        "ssgnp_0p450v_m40c_cworst_CCworst_T",
        "derived from the conduction structure, not node naming "
        "(proven by the rename-invariance gate)"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(18 if i == 0 else 15)
        r.font.color.rgb = INK if i == 0 else MUTE
        r.font.name = "Arial"

    # ---- content slides ----
    area_w, area_h = Inches(12.6), Inches(5.7)
    area_left, area_top = Inches(0.366), Inches(1.55)
    for title, subtitle, fig in SLIDES:
        slide = prs.slides.add_slide(blank)
        _title_box(slide, title, Inches(0.28), 23, INK)
        sb = _title_box(slide, subtitle, Inches(0.95), 13, MUTE, bold=False)
        png = raster(fig)
        iw, ih = Image.open(png).size
        w, h = _fit(iw, ih, area_w, area_h)
        left = area_left + (area_w - w) // 2
        top = area_top + (area_h - h) // 2
        slide.shapes.add_picture(png, left, top, width=w, height=h)

    prs.save(OUT)
    print("wrote", OUT, "(%d slides)" % len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
